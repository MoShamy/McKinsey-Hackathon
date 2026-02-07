import os
import json
import re
from pptx import Presentation
from openai import OpenAI

# --------------------------------------------------
# 1. OpenAI client setup
# --------------------------------------------------
# Set this in your terminal BEFORE running:
# export OPENAI_API_KEY="your_real_api_key"

client = OpenAI(
    api_key=os.getenv("OPENAI_API_KEY"),
    base_url="https://openai.prod.ai-gateway.quantumblack.com/2907fa1c-683f-4cf7-8555-8c34a16c88d8/v1"
)

# --------------------------------------------------
# 2. Dot-dash outline
# --------------------------------------------------
dot_dash_outline = """
- Introduction
  - Problem Statement
  - Objective
- Methodology
  - Data Collection
  - Analysis
- Results
  - Findings
  - Interpretation
- Conclusion
"""

# --------------------------------------------------
# 3. Generate slide content using GPT
# --------------------------------------------------
def generate_slide_content(outline: str):
    prompt = f"""
You are an AI Presentation Designer.

Input: A content outline in dot-dash format:
{outline}

Task:
1. Create a slide for each main point.
2. Convert sub-bullets into concise slide bullet points.
3. Add presenter notes for each slide.
4. Specify slide type: Title, Bullet, Chart, Diagram, Image.
5. Output ONLY valid JSON as an array with:
   - slide_title
   - slide_type
   - content (list of bullets)
   - presenter_notes
"""

    response = client.chat.completions.create(
        model="gpt-4",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3
    )

    text_output = response.choices[0].message.content.strip()

    # --------------------------------------------------
    # Safely extract JSON
    # --------------------------------------------------
    try:
        match = re.search(r"\[.*\]", text_output, re.DOTALL)
        if not match:
            raise ValueError("No JSON array found")

        return json.loads(match.group(0))

    except Exception as e:
        print("Failed to parse JSON:", e)
        print("Raw model output:\n", text_output)
        return None

# --------------------------------------------------
# 4. Create PowerPoint presentation
# --------------------------------------------------
def create_presentation(slide_data, output_file="presentation.pptx"):
    prs = Presentation()

    for slide in slide_data:
        slide_type = slide.get("slide_type", "").lower()

        # Layout selection
        layout_index = 0 if slide_type == "title" else 1
        layout_index = min(layout_index, len(prs.slide_layouts) - 1)

        s = prs.slides.add_slide(prs.slide_layouts[layout_index])

        # Title
        s.shapes.title.text = slide.get("slide_title", "")

        # Content bullets
        if layout_index == 1 and slide.get("content"):
            tf = s.placeholders[1].text_frame
            tf.clear()

            for i, bullet in enumerate(slide["content"]):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = bullet
                p.level = 0

        # Presenter notes
        if slide.get("presenter_notes"):
            s.notes_slide.notes_text_frame.text = slide["presenter_notes"]

    prs.save(output_file)
    print(f"Presentation saved as: {output_file}")

# --------------------------------------------------
# 5. Run pipeline
# --------------------------------------------------
slides = generate_slide_content(dot_dash_outline)

if slides:
    create_presentation(slides)
else:
    print("Slide generation failed.")