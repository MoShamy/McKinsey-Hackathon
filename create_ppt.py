import copy
from io import BytesIO
import math
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- MCKINSEY-STYLE PALETTE ---
NAVY_BG = RGBColor(15, 23, 42)       # Dark Navy
WHITE_TEXT = RGBColor(255, 255, 255) # White
ACCENT_BLUE = RGBColor(56, 189, 248) # Cyan/Light Blue Accent
DARK_TEXT = RGBColor(30, 41, 59)     # Slate 800 for body text

def fit_font_size(text, max_len=50, default_size=32, min_size=18):
    """
    Smart Helper: Heuristic to reduce font size if text is too long.
    """
    if not text: return default_size
    length = len(text)
    
    if length < max_len:
        return default_size
    elif length < max_len * 1.5:
        return max(min_size, default_size - 4)
    elif length < max_len * 2:
        return max(min_size, default_size - 8)
    else:
        return min_size

def add_header(slide, title_text):
    """Draws the Navy Header bar and places the title safely."""
    # 1. Navy Bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY_BG
    header.line.fill.background() # No border

    # 2. Title Text (Auto-Scaling)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Arial"
    p.font.bold = True
    p.font.color.rgb = WHITE_TEXT
    
    # Smart Size: If title is huge, shrink it.
    p.font.size = Pt(fit_font_size(title_text, max_len=40, default_size=32, min_size=24))

def create_title_slide(prs, slide_data):
    """Creates a sleek, dark title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    
    title_text = slide_data.get("title", "Executive Presentation")
    
    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = NAVY_BG

    # Title
    t_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = t_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.bold = True
    p.font.size = Pt(fit_font_size(title_text, max_len=30, default_size=48, min_size=36))
    p.font.color.rgb = WHITE_TEXT

    # Accent Line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.5), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # Subtitle (Optional)
    subtitle_text = "Strategic Overview"
    if subtitle_text:
        s_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
        sf = s_box.text_frame
        p2 = sf.paragraphs[0]
        p2.text = subtitle_text
        p2.alignment = PP_ALIGN.CENTER
        p2.font.name = "Arial"
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(200, 200, 200)

    # --- FIX: FORCE NOTES CREATION ---
    notes_text = slide_data.get("speaker_notes", "")
    if notes_text:
        notes_slide = slide.notes_slide # This creates the notes slide if missing
        notes_slide.notes_text_frame.text = notes_text

def create_smart_content_slide(prs, slide_data):
    """
    Decides layout based on content density.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    
    title = slide_data.get("title", "Untitled")
    bullets = slide_data.get("bullets", [])
    
    # 1. Draw Header
    add_header(slide, title)

    # 2. Layout Logic
    total_chars = sum(len(b) for b in bullets)
    num_bullets = len(bullets)
    
    use_two_columns = False
    if num_bullets > 5 or total_chars > 400:
        use_two_columns = True

    # 3. Render Body
    if not use_two_columns:
        # --- SINGLE COLUMN LAYOUT ---
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        tf = box.text_frame
        tf.word_wrap = True
        
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0
            p.space_after = Pt(14)
            p.font.name = "Arial"
            p.font.color.rgb = DARK_TEXT
            p.font.size = Pt(fit_font_size(b, max_len=80, default_size=20, min_size=14))
            
    else:
        # --- TWO COLUMN LAYOUT ---
        mid_point = math.ceil(num_bullets / 2)
        left_bullets = bullets[:mid_point]
        right_bullets = bullets[mid_point:]
        
        # Left Box
        box1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.25), Inches(5.5))
        tf1 = box1.text_frame
        tf1.word_wrap = True
        for b in left_bullets:
            p = tf1.add_paragraph()
            p.text = b
            p.level = 0
            p.space_after = Pt(12)
            p.font.name = "Arial"
            p.font.color.rgb = DARK_TEXT
            p.font.size = Pt(16)

        # Right Box
        box2 = slide.shapes.add_textbox(Inches(5.25), Inches(1.5), Inches(4.25), Inches(5.5))
        tf2 = box2.text_frame
        tf2.word_wrap = True
        for b in right_bullets:
            p = tf2.add_paragraph()
            p.text = b
            p.level = 0
            p.space_after = Pt(12)
            p.font.name = "Arial"
            p.font.color.rgb = DARK_TEXT
            p.font.size = Pt(16)

    # --- FIX: FORCE NOTES CREATION ---
    notes_text = slide_data.get("speaker_notes", "")
    if notes_text:
        notes_slide = slide.notes_slide # This creates the notes slide if missing
        notes_slide.notes_text_frame.text = notes_text

def generate_pptx(json_data, template_file=None):
    # Setup
    prs = Presentation() 
    
    slides_data = json_data.get("slides", [])
    if not slides_data:
        # Fallback if no slides found
        return BytesIO()

    # 1. Title Slide (Assumes first slide in JSON is title)
    create_title_slide(prs, slides_data[0])

    # 2. Content Slides (Rest of the list)
    if len(slides_data) > 1:
        for i in range(1, len(slides_data)):
            create_smart_content_slide(prs, slides_data[i])

    # Export
    pptx_stream = BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    return pptx_stream